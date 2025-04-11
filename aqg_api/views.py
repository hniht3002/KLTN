from django.shortcuts import render
import os
from django.conf import settings
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
from google import genai
from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework import status
from rest_framework.parsers import MultiPartParser, FormParser
from .models import Document
from .serializers import DocumentUploadSerializer
from dotenv import load_dotenv
from langchain_community.vectorstores import FAISS
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from aqg_api.aqg.main import AQG

load_dotenv()
api_key = os.getenv("GOOGLE_API_KEY")
        
if not api_key:
    raise Exception("GOOGLE_API_KEY not found in .env file")

# Configure the Gemini API
client = genai.Client(api_key = api_key)

class UploadFileView(APIView):
    parser_classes = (MultiPartParser, FormParser)
    
    def post(self, request, *args, **kwargs):
        serializer = DocumentUploadSerializer(data=request.data, context={'request': request})
        
        content = ""
        if serializer.is_valid():
            # Save the document first
            document = serializer.save()
            
            # Extract text from PDF
            if document.file_type == 'pdf':
                try:
                    # Open the PDF file
                    pdf_path = os.path.join(settings.MEDIA_ROOT, str(document.file))
                    with open(pdf_path, 'rb') as pdf_file:
                        # Create PDF reader object
                        pdf_reader = PyPDF2.PdfReader(pdf_file)

                        for page in pdf_reader.pages:
                            content += page.extract_text()
                        
                        # Save extracted text to document
                        document.content = content
                        document.save()
                except Exception as e:
                    # Delete document if text extraction fails
                    document.delete()
                    return Response({
                        'error': f'Failed to extract text from PDF: {str(e)}'
                    }, status=status.HTTP_400_BAD_REQUEST)
                    
            elif document.file_type == 'docx':
                try:
                    
                    # Open the Word file
                    docx_path = os.path.join(settings.MEDIA_ROOT, str(document.file))
                    doc = DocxDocument(docx_path)
                    
                    for para in doc.paragraphs:
                        content += para.text + "\n"
                
                    
                    # Save extracted text to document
                    document.content = content
                    document.save()
                
                except Exception as e:
                    # Delete document if text extraction fails
                    document.delete()
                    return Response({
                        'error': f'Failed to extract text from Word document: {str(e)}'
                    }, status=status.HTTP_400_BAD_REQUEST)
                    
            elif document.file_type == 'pptx':
                try:                    
                    # Open the PowerPoint file
                    pptx_path = os.path.join(settings.MEDIA_ROOT, str(document.file))
                    prs = Presentation(pptx_path)

                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                content += shape.text + "\n"
                    
                    # Save extracted text to document
                    document.content = content
                    document.save()
                
                except Exception as e:
                    # Delete document if text extraction fails 
                    document.delete()
                    return Response({
                        'error': f'Failed to extract text from PowerPoint: {str(e)}'
                    }, status=status.HTTP_400_BAD_REQUEST)
            
            store_document_embeddings(content)

            return Response(serializer.data, status=status.HTTP_201_CREATED)
            
def embed_content(content):
    from langchain.text_splitter import RecursiveCharacterTextSplitter
    
    # Create text splitter with target chunk size of ~500 words
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=2000,  # Approx 500 words
        chunk_overlap=200,  # Approx 50 words
        length_function=len,
        separators=["\n\n", "\n", " ", ""]
    )
    
    # Split text into chunks
    chunks = text_splitter.split_text(content)
    print(f"Content split into {len(chunks)} chunks")
    
    # Generate embeddings for all chunks using Gemini
    embedded_chunks = []

    for chunk in chunks:
        embedding = client.models.embed_content(
            model="text-embedding-004",
            contents=chunk
        ).embeddings[0].values
        embedded_chunks.append(embedding)
    
    return embedded_chunks

def store_document_embeddings(content):
    # Create text splitter
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=2000,
        chunk_overlap=200,
        length_function=len,
        separators=["\n\n", "\n", " ", ""]
    )
    
    # Split text into chunks
    chunks = text_splitter.split_text(content)
    print(f"Content split into {len(chunks)} chunks")
    
    # Initialize embeddings
    embedding = GoogleGenerativeAIEmbeddings(
        google_api_key=os.getenv("GOOGLE_API_KEY"),
        model="models/text-embedding-004")
    
    # Create FAISS index from texts
    vectorstore = FAISS.from_texts(chunks, embedding)
 
    vectorstore.save_local(settings.VECTOR_DB_DIR)
    
    return vectorstore

class GenerateQuestionsView(APIView):
    
    def post(self, request, *args, **kwargs):
        try:
            # Get topic from request
            topic = request.data.get('topic')
            num_questions = int(request.data.get('num_questions', 3))
            
            if not topic:
                return Response(
                    {'error': 'topic is required'}, 
                    status=status.HTTP_400_BAD_REQUEST
                )
                
            aqg = AQG(topic)
            
            relevant_documents = aqg.retriever(topic = topic, num_documents = 3)
                        
            # Combine relevant chunks into context
            context = "\n\n".join([doc.page_content for doc, _ in relevant_documents])
            
            response = aqg.generator(num_questions, context)
            
            return Response(response, status=status.HTTP_200_OK)

        except Exception as e:
            return Response(
                {'error': str(e)}, 
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )

class UploadAndGenerateView(APIView):
    parser_classes = (MultiPartParser, FormParser)
    def post(self, request, *args, **kwargs):
        serializer = DocumentUploadSerializer(data=request.data, context={'request': request})
        
        try:
            content = ""
            if serializer.is_valid():
                # Save the document first
                document = serializer.save()           
                # Extract text from PDF
                if document.file_type == 'pdf':
                    try:
                        # Open the PDF file
                        pdf_path = os.path.join(settings.MEDIA_ROOT, str(document.file))
                        with open(pdf_path, 'rb') as pdf_file:
                            # Create PDF reader object
                            pdf_reader = PyPDF2.PdfReader(pdf_file)

                            for page in pdf_reader.pages:
                                content += page.extract_text()
                            
                            # Save extracted text to document
                            document.content = content
                            document.save()
                    except Exception as e:
                        # Delete document if text extraction fails
                        document.delete()
                        return Response({
                            'error': f'Failed to extract text from PDF: {str(e)}'
                        }, status=status.HTTP_400_BAD_REQUEST)
                        
                elif document.file_type == 'docx':
                    try:
                        
                        # Open the Word file
                        docx_path = os.path.join(settings.MEDIA_ROOT, str(document.file))
                        doc = DocxDocument(docx_path)
                        
                        for para in doc.paragraphs:
                            content += para.text + "\n"
                    
                        
                        # Save extracted text to document
                        document.content = content
                        document.save()
                    
                    except Exception as e:
                        # Delete document if text extraction fails
                        document.delete()
                        return Response({
                            'error': f'Failed to extract text from Word document: {str(e)}'
                        }, status=status.HTTP_400_BAD_REQUEST)
                        
                elif document.file_type == 'pptx':
                    try:                    
                        # Open the PowerPoint file
                        pptx_path = os.path.join(settings.MEDIA_ROOT, str(document.file))
                        prs = Presentation(pptx_path)

                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    content += shape.text + "\n"
                        
                        # Save extracted text to document
                        document.content = content
                        document.save()
                    
                    except Exception as e:
                        # Delete document if text extraction fails 
                        document.delete()
                        return Response({
                            'error': f'Failed to extract text from PowerPoint: {str(e)}'
                        }, status=status.HTTP_400_BAD_REQUEST)
                
                store_document_embeddings(content)
            
            aqg = AQG("")
            
            response = aqg.generator(num_questions = 10, context = content)
            
            return Response(response, status=status.HTTP_200_OK)
        
        except Exception as e:
            return Response(
                {'error': str(e)}, 
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )

class EvaluateQuestionsView(APIView):
    def post(self, request, *args, **kwargs):
        try:
            aqg = AQG(topic="")
            
            type = request.data.get('type')
            if type == "generator":
                evaluation_result = aqg.generator_evaluate()
            else:
                evaluation_result = aqg.generate_evaluate_response()
            
            return Response(evaluation_result, status=status.HTTP_200_OK)
            
        except Exception as e:
            return Response(
                {'error': str(e)}, 
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )